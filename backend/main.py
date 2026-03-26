"""
FastAPI backend for PowerPoint Narration Generator.

POST /api/parse   — parse Word doc, return slide list (no TTS yet)
POST /api/process — full pipeline: parse Word → TTS → embed audio → return PPTX
GET  /            — serve the wizard UI
"""
import asyncio
import base64
import io
import json
import queue as _queue
import tempfile
import threading
import os
import zipfile
from pathlib import Path

from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.responses import StreamingResponse, HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware

from word_parser import extract_slides
from pptx_script_parser import extract_slides_from_pptx
from tts_client import synthesize_to_mp3
from pptx_builder import embed_audio_into_pptx
from translator import translate_for_voice
from quality_checker import run_quality_check
from ai_pptx_generator import build_ai_presentation
from video_exporter import export_video
from pptx import Presentation

app = FastAPI(title="PowerPoint Narration Generator")


_OLD_DOC_MAGIC = b"\xd0\xcf\x11\xe0"   # OLE2 / legacy .doc signature
_ECMA_MAGIC    = b"PK\x03\x04"          # ZIP / modern Office Open XML signature


def _parse_script(filename: str, file_bytes: bytes) -> list[dict]:
    """Auto-detect script format (docx or pptx) and extract slides."""
    try:
        if filename.lower().endswith(".pptx"):
            return extract_slides_from_pptx(file_bytes)
        return extract_slides(file_bytes)
    except zipfile.BadZipFile:
        # Give a more specific hint based on the file's magic bytes
        if file_bytes[:4] == _OLD_DOC_MAGIC:
            detail = (
                f"'{filename}' appears to be a legacy Word 97-2003 (.doc) file saved "
                f"with a .docx extension. Please open it in Word, choose "
                f"File → Save As → Word Document (.docx), and re-upload."
            )
        elif file_bytes[:4] != _ECMA_MAGIC:
            detail = (
                f"'{filename}' could not be opened — it may be password-protected or "
                f"corrupted. If the document has a password, remove it in Word "
                f"(File → Info → Protect Document → Encrypt with Password → clear the "
                f"password) and re-upload."
            )
        else:
            detail = (
                f"'{filename}' could not be read as a valid "
                f"{'PPTX' if filename.lower().endswith('.pptx') else 'DOCX'} file. "
                f"Please check the file is not corrupted and matches its extension."
            )
        raise HTTPException(status_code=422, detail=detail)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Failed to parse script: {exc}")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# Serve static frontend files
_frontend = Path(__file__).parent.parent / "frontend"
if _frontend.exists():
    app.mount("/static", StaticFiles(directory=str(_frontend)), name="static")


_ENV_BANNER_MSG = os.environ.get("APP_BANNER_MESSAGE", "").strip()
_ENV_BANNER_HTML = (
    f'<div class="env-banner" role="status">{_ENV_BANNER_MSG}</div>'
    if _ENV_BANNER_MSG else ""
)


@app.get("/", response_class=HTMLResponse)
async def index():
    index_file = _frontend / "index.html"
    if index_file.exists():
        html = index_file.read_text(encoding="utf-8")
        html = html.replace("{{ENV_BANNER}}", _ENV_BANNER_HTML)
        return HTMLResponse(html)
    return HTMLResponse("<h1>Frontend not found</h1>", status_code=404)


@app.post("/api/parse")
async def parse_doc(
    script: UploadFile = File(...),
    pptx: UploadFile = File(None),
    ai_mode: str = Form("false"),
):
    """
    Step 1 of the wizard: parse the script and return slide list.
    When ai_mode=true, pptx is optional — slide count comes from the script alone.
    """
    script_bytes = await script.read()

    slides = _parse_script(script.filename or "", script_bytes)

    if ai_mode.lower() == "true" or pptx is None or pptx.filename == "":
        return JSONResponse({
            "slides": slides,
            "pptx_slide_count": len(slides),
            "word_slide_count": len(slides),
            "ai_mode": True,
        })

    pptx_bytes = await pptx.read()
    prs = Presentation(io.BytesIO(pptx_bytes))
    pptx_slide_count = len(prs.slides)

    return JSONResponse({
        "slides": slides,
        "pptx_slide_count": pptx_slide_count,
        "word_slide_count": len(slides),
        "ai_mode": False,
    })


@app.post("/api/process")
async def process(
    script: UploadFile = File(...),
    pptx: UploadFile = File(...),
    voice: str = Form("en-US-JennyNeural"),
    slide_mapping: str = Form("{}"),  # JSON: {"0": 0, "1": 1, ...} word→pptx index
):
    """
    Full pipeline:
      1. Parse Word doc into slide blocks.
      2. Synthesize each block to MP3 via Azure TTS.
      3. Embed audio into PPTX slides per the mapping.
      4. Return the modified PPTX as a download.
    """
    script_bytes = await script.read()
    pptx_bytes = await pptx.read()

    # Parse script (Word or PowerPoint)
    slides = _parse_script(script.filename or "", script_bytes)

    prs = Presentation(io.BytesIO(pptx_bytes))
    pptx_slide_count = len(prs.slides)

    # Parse mapping: word slide index → pptx slide index
    try:
        mapping: dict[str, int] = json.loads(slide_mapping)
    except Exception:
        mapping = {}

    # Default mapping: pair by position up to min(word, pptx)
    if not mapping:
        for i in range(min(len(slides), pptx_slide_count)):
            mapping[str(i)] = i

    # Build per-pptx-slide audio list
    slide_audio: list[bytes | None] = [None] * pptx_slide_count

    for word_idx_str, pptx_idx in mapping.items():
        word_idx = int(word_idx_str)
        if word_idx >= len(slides):
            continue
        if pptx_idx >= pptx_slide_count:
            continue

        text = slides[word_idx]["text"]
        if not text.strip():
            continue

        print(f"[Process] Synthesizing slide {word_idx + 1} → PPTX slide {pptx_idx + 1} ({len(text)} chars)", flush=True)
        try:
            translated = translate_for_voice(text, voice=voice)
            audio = synthesize_to_mp3(translated, voice=voice)
            slide_audio[pptx_idx] = audio if audio else None
        except Exception as exc:
            raise HTTPException(
                status_code=502,
                detail=f"TTS failed for slide {word_idx + 1}: {exc}",
            )

    # Embed audio
    print(f"[Process] Embedding audio into PPTX ({sum(1 for a in slide_audio if a)} slides with audio)...", flush=True)
    result_bytes = embed_audio_into_pptx(pptx_bytes, slide_audio)
    print(f"[Process] Done — {len(result_bytes):,} bytes output", flush=True)

    return StreamingResponse(
        io.BytesIO(result_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=narrated_presentation.pptx"},
    )


@app.post("/api/generate-ai")
async def generate_ai(
    script: UploadFile = File(...),
    voice: str = Form("en-US-JennyNeural"),
):
    """
    AI generation pipeline — streams newline-delimited JSON progress events:
      {"type":"progress","slide":N,"total":M,"phase":"structure|image|build|tts","message":"..."}
    Ends with:
      {"type":"done","pptx":"<base64 encoded PPTX>"}
    or on failure:
      {"type":"error","message":"..."}
    """
    script_bytes = await script.read()
    slides = _parse_script(script.filename or "", script_bytes)
    total = len(slides)

    print(f"[AI-Gen] Starting AI generation: {total} slides, voice={voice}", flush=True)

    progress_q: _queue.Queue = _queue.Queue()

    def _run():
        try:
            _phase_labels = {
                "structure": "Structuring slide {n} of {t} with GPT…",
                "image":     "Generating image for slide {n} of {t}…",
                "build":     "Building slide {n} of {t}…",
                "tts":       "Synthesising audio for slide {n} of {t}…",
            }

            def on_progress(slide_num: int, tot: int, phase: str):
                msg = _phase_labels.get(phase, f"Processing slide {slide_num} of {tot}…")
                msg = msg.format(n=slide_num, t=tot)
                progress_q.put({
                    "type": "progress",
                    "slide": slide_num,
                    "total": tot,
                    "phase": phase,
                    "message": msg,
                })

            # Steps 1–3: GPT structure + images + PPTX build
            pptx_bytes = build_ai_presentation(slides, on_progress=on_progress)

            # Step 4: TTS synthesis
            slide_audio: list[bytes | None] = []
            for idx, slide in enumerate(slides):
                text = slide.get("text", "").strip()
                if not text:
                    slide_audio.append(None)
                    continue
                on_progress(idx + 1, total, "tts")
                try:
                    translated = translate_for_voice(text, voice=voice)
                    audio = synthesize_to_mp3(translated, voice=voice)
                    slide_audio.append(audio if audio else None)
                except Exception as exc:
                    print(f"[AI-Gen] TTS failed for slide {idx + 1}: {exc}", flush=True)
                    slide_audio.append(None)

            # Step 5: embed audio + encode result
            result_bytes = embed_audio_into_pptx(pptx_bytes, slide_audio)
            print(f"[AI-Gen] Done — {len(result_bytes):,} bytes", flush=True)
            progress_q.put({"type": "done", "pptx": base64.b64encode(result_bytes).decode()})

        except Exception as exc:
            progress_q.put({"type": "error", "message": str(exc)})

    thread = threading.Thread(target=_run, daemon=True)
    thread.start()

    async def _event_stream():
        while True:
            try:
                event = progress_q.get_nowait()
                yield json.dumps(event) + "\n"
                if event["type"] in ("done", "error"):
                    break
            except _queue.Empty:
                await asyncio.sleep(0.1)

    return StreamingResponse(_event_stream(), media_type="application/x-ndjson")


@app.post("/api/quality-check")
async def quality_check(
    script: UploadFile = File(...),
    pptx: UploadFile = File(...),
    voice: str = Form("en-US-JennyNeural"),
):
    """
    Quality Check agent:
      1. Parse original Word script for slide texts.
      2. Extract MP3s from the narrated PPTX and transcribe each via Azure STT.
      3. Compare each transcription against the script section using GPT-5.2.
      4. Return per-slide confidence scores and issues.
    """
    script_bytes = await script.read()
    pptx_bytes = await pptx.read()

    slides = _parse_script(script.filename or "", script_bytes)
    print(f"[QA] Starting quality check: {len(slides)} script slides, voice={voice}", flush=True)

    try:
        results = run_quality_check(pptx_bytes, slides, voice=voice)
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Quality check failed: {exc}")

    return JSONResponse({"results": results})


@app.post("/api/export-video")
async def export_video_route(
    pptx: UploadFile = File(...),
):
    """
    Convert a narrated PPTX to an MP4 slideshow.
    Streams newline-delimited JSON progress events, ending with:
      {"type":"done","mp4":"<base64>"}
    """
    pptx_bytes = await pptx.read()
    progress_q: _queue.Queue = _queue.Queue()

    def _run():
        try:
            _phase_labels = {
                "export": "Exporting slides as images…",
                "encode": "Encoding slide {n} of {t}…",
                "concat": "Combining clips into video…",
            }

            def on_progress(slide_num: int, tot: int, phase: str):
                msg = _phase_labels.get(phase, f"Processing slide {slide_num} of {tot}…")
                msg = msg.format(n=slide_num, t=tot)
                progress_q.put({
                    "type": "progress",
                    "slide": slide_num,
                    "total": tot,
                    "phase": phase,
                    "message": msg,
                })

            mp4_bytes = export_video(pptx_bytes, on_progress=on_progress)
            progress_q.put({"type": "done", "mp4": base64.b64encode(mp4_bytes).decode()})
        except Exception as exc:
            progress_q.put({"type": "error", "message": str(exc)})

    thread = threading.Thread(target=_run, daemon=True)
    thread.start()

    async def _event_stream():
        while True:
            try:
                event = progress_q.get_nowait()
                yield json.dumps(event) + "\n"
                if event["type"] in ("done", "error"):
                    break
            except _queue.Empty:
                await asyncio.sleep(0.1)

    return StreamingResponse(_event_stream(), media_type="application/x-ndjson")

