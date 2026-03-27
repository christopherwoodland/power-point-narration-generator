"""
AI-powered PowerPoint generator.

For each slide in the parsed script:
  1. Call GPT-5.2 to produce structured JSON (title, bullets, image_prompt, include_image).
  2. For slides where include_image=true, call gpt-image-1.5 to generate an image.
  3. Build a clean PPTX from scratch using python-pptx.
  4. Return the PPTX as bytes (audio is added by the caller via pptx_builder).

Uses DefaultAzureCredential — no API keys.
Endpoint: https://bhs-development-public-foundry-r.cognitiveservices.azure.com
Deployments: gpt-5.2 (chat), gpt-image-1.5 (images)
"""
import base64
import io
import json
import os
import re

import httpx
from azure_credential import get_credential
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

COGNITIVE_SERVICES_SCOPE = "https://cognitiveservices.azure.com/.default"

_ENDPOINT = os.environ.get(
    "AZURE_OPENAI_ENDPOINT",
    "https://bhs-development-public-foundry-r.cognitiveservices.azure.com",
).rstrip("/")

_CHAT_DEPLOYMENT = os.environ.get("AZURE_OPENAI_DEPLOYMENT", "gpt-5.2")
_IMG_DEPLOYMENT  = os.environ.get("AZURE_IMAGE_DEPLOYMENT", "gpt-image-1.5")
_CHAT_API_VER    = "2025-01-01-preview"
_IMG_API_VER     = "2024-02-01"

_CHAT_URL = (
    f"{_ENDPOINT}/openai/deployments/{_CHAT_DEPLOYMENT}"
    f"/chat/completions?api-version={_CHAT_API_VER}"
)
_IMG_URL = (
    f"{_ENDPOINT}/openai/deployments/{_IMG_DEPLOYMENT}"
    f"/images/generations?api-version={_IMG_API_VER}"
)

# ── Slide dimensions (widescreen 13.33" × 7.5") ───────────────────────────
_W = Inches(13.33)
_H = Inches(7.5)

# Theme colours
_DARK_BLUE = RGBColor(0x00, 0x3A, 0x6E)
_MID_BLUE  = RGBColor(0x00, 0x78, 0xD4)
_DARK_GREY = RGBColor(0x20, 0x1F, 0x1E)
_MID_GREY  = RGBColor(0x60, 0x5E, 0x5C)
_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_BG  = RGBColor(0xF8, 0xF9, 0xFA)

# ── GPT system prompt ──────────────────────────────────────────────────────
_STRUCTURE_SYSTEM = """You are a professional presentation designer.
Given a narration script for one presentation slide, return a JSON object describing how to
visually design that slide. Respond ONLY with valid JSON — no markdown fences, no explanation.

Schema:
{
  "title": "<slide title, max 12 words>",
  "bullets": ["<point 1>", "<point 2>", "<point 3>"],
  "include_image": <true or false>,
  "image_prompt": "<DALL-E image generation prompt if include_image is true, else empty string>"
}

Rules:
- bullets: 2–4 concise bullet points summarising the key message. Never repeat the title.
- include_image: true for content slides; false for purely textual/agenda/closing slides.
- image_prompt: a vivid, specific prompt for a professional business/educational illustration
  that complements the slide content. Do NOT mention text, words, captions, or logos in the prompt.
  Request a clean, photorealistic or flat-design illustration on a light background.
"""


def _get_aad_token() -> str:
    return get_credential().get_token(COGNITIVE_SERVICES_SCOPE).token


def _structure_slide(narration_text: str, slide_num: int) -> dict:
    """Ask GPT-5.2 to return structured slide content for a single slide."""
    token = _get_aad_token()
    payload = {
        "messages": [
            {"role": "system", "content": _STRUCTURE_SYSTEM},
            {"role": "user", "content": f"Slide {slide_num} narration:\n\n{narration_text}"},
        ],
        "temperature": 0.4,
        "max_completion_tokens": 600,
        "response_format": {"type": "json_object"},
    }
    resp = httpx.post(
        _CHAT_URL,
        headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"},
        content=json.dumps(payload),
        timeout=60,
    )
    if not resp.is_success:
        print(f"[AI-Gen] GPT error {resp.status_code}: {resp.text[:400]}", flush=True)
    resp.raise_for_status()
    content = resp.json()["choices"][0]["message"]["content"]
    return json.loads(content)


def _generate_image(prompt: str, slide_num: int) -> bytes | None:
    """Call gpt-image-1.5 and return PNG bytes, or None on failure."""
    try:
        token = _get_aad_token()
        payload = {
            "prompt": prompt,
            "n": 1,
            "size": "1024x1024",
        }
        resp = httpx.post(
            _IMG_URL,
            headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"},
            content=json.dumps(payload),
            timeout=120,
        )
        if not resp.is_success:
            print(f"[AI-Gen] Image error slide {slide_num} {resp.status_code}: {resp.text[:300]}", flush=True)
            return None
        b64 = resp.json()["data"][0]["b64_json"]
        return base64.b64decode(b64)
    except Exception as exc:
        print(f"[AI-Gen] Image generation failed for slide {slide_num}: {exc}", flush=True)
        return None


# ── python-pptx helpers ────────────────────────────────────────────────────

def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H
    return prs


def _blank_slide(prs: Presentation):
    """Add a truly blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]  # index 6 = Blank in all default templates
    return prs.slides.add_slide(blank_layout)


def _set_slide_bg(slide, colour: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _add_title_bar(slide, title_text: str):
    """Dark blue bar at top containing slide title."""
    bar_h = Inches(1.1)
    txBox = slide.shapes.add_textbox(Inches(0), Inches(0), _W, bar_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Blue background — use the python-pptx fill API to avoid duplicate noFill+solidFill
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = _DARK_BLUE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.bold   = True
    run.font.size   = Pt(26)
    run.font.color.rgb = _WHITE
    # Left padding
    from pptx.oxml.ns import qn as _qn
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(0.3))))
    pPr.set("indent", "0")
    tf.margin_top    = Inches(0.18)
    tf.margin_bottom = Inches(0.0)
    tf.margin_left   = Inches(0.3)


def _add_bullets(slide, bullets: list[str], x, y, w, h):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.15)
    tf.margin_top    = Inches(0.1)

    for i, bullet in enumerate(bullets):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        from pptx.oxml.ns import qn
        pPr = para._p.get_or_add_pPr()
        pPr.set("marL", str(int(Inches(0.2))))
        pPr.set("indent", str(int(Inches(-0.2))))
        pPr.set("spc", "0")

        run = para.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(17)
        run.font.color.rgb = _DARK_GREY
        # Space before each bullet after the first
        if i > 0:
            from lxml import etree
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
            spcPts.set("val", "180")  # 18pt


def _add_image_to_slide(slide, image_bytes: bytes, x, y, w, h):
    img_stream = io.BytesIO(image_bytes)
    slide.shapes.add_picture(img_stream, x, y, w, h)


def _add_narration_label(slide):
    """Small 'AI Generated' watermark bottom-right."""
    txBox = slide.shapes.add_textbox(
        _W - Inches(2.4), _H - Inches(0.35),
        Inches(2.3), Inches(0.3),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "✦ AI Generated Presentation"
    run.font.size = Pt(8)
    run.font.color.rgb = _MID_GREY
    run.font.italic = True


# ── Main public API ────────────────────────────────────────────────────────

def build_ai_presentation(
    slides: list[dict],
    on_progress=None,  # callable(slide_num, total, phase: str) | None
) -> bytes:
    """
    Given a list of {"title": ..., "text": ...} dicts (from the script parser),
    generate a full PPTX using GPT-5.2 for structure and gpt-image-1.5 for images.
    Returns raw PPTX bytes (no audio — caller adds audio via pptx_builder).

    on_progress(slide_num, total, phase) is called before each major phase so
    the caller can stream progress updates. phase is one of: "structure", "image", "build".
    """
    total = len(slides)
    prs = _new_presentation()

    def _progress(slide_num: int, phase: str):
        if on_progress is not None:
            try:
                on_progress(slide_num, total, phase)
            except Exception:
                pass

    for idx, slide_data in enumerate(slides):
        slide_num = idx + 1
        narration = slide_data.get("text", "").strip()
        fallback_title = slide_data.get("title", f"Slide {slide_num}")

        # ── Step 1: ask GPT to structure the slide ─────────────────────────
        print(f"[AI-Gen] Structuring slide {slide_num}/{total}...", flush=True)
        _progress(slide_num, "structure")
        try:
            structured = _structure_slide(narration or fallback_title, slide_num)
        except Exception as exc:
            print(f"[AI-Gen] GPT structure failed for slide {slide_num}: {exc}", flush=True)
            structured = {
                "title": fallback_title,
                "bullets": [narration[:120]] if narration else [],
                "include_image": False,
                "image_prompt": "",
            }

        title_text   = structured.get("title") or fallback_title
        bullets      = structured.get("bullets") or []
        include_img  = structured.get("include_image", False)
        image_prompt = structured.get("image_prompt", "")

        # ── Step 2: generate image ─────────────────────────────────────────
        image_bytes: bytes | None = None
        if include_img and image_prompt:
            print(f"[AI-Gen] Generating image for slide {slide_num}...", flush=True)
            _progress(slide_num, "image")
            image_bytes = _generate_image(image_prompt, slide_num)

        # ── Step 3: build the slide ────────────────────────────────────────
        _progress(slide_num, "build")
        sl = _blank_slide(prs)
        _set_slide_bg(sl, _LIGHT_BG)
        _add_title_bar(sl, title_text)

        content_top  = Inches(1.25)
        content_h    = _H - Inches(1.6)
        pad          = Inches(0.3)

        if image_bytes:
            # Left: bullets  |  Right: image
            text_w = Inches(6.2)
            img_x  = Inches(6.8)
            img_w  = _W - img_x - pad
            img_h  = min(content_h, img_w)          # square-ish, constrained to height
            img_y  = content_top + (content_h - img_h) // 2

            _add_bullets(sl, bullets, pad, content_top, text_w, content_h)
            _add_image_to_slide(sl, image_bytes, img_x, img_y, img_w, img_h)
        else:
            # No image — text fills full width
            _add_bullets(sl, bullets, pad, content_top, _W - 2 * pad, content_h)

        _add_narration_label(sl)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
