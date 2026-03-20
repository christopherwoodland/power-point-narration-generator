"""
Embed per-slide MP3 audio into a PowerPoint file.

Uses python-pptx's OPC Part API for proper package handling:
  1. Adds MP3 and icon PNG as proper OPC Parts via slide.part.relate_to().
  2. Inserts a <p:pic> element that references the audio via r:link.
  3. Inserts <p:timing> so the audio auto-plays on slide entry.
"""
import io
import struct as _struct
import zlib as _zlib
from lxml import etree
from pptx import Presentation
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.oxml.ns import qn


# OPC relationship types
_AUDIO_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio"
)
# Modern PowerPoint also requires a second rel with the 2007 media type pointing to the same file
_MEDIA_REL_TYPE = (
    "http://schemas.microsoft.com/office/2007/relationships/media"
)
_IMAGE_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
)


def _make_tiny_png() -> bytes:
    """Generate a minimal 1x1 white PNG for use as the audio icon placeholder."""
    def _chunk(tag: bytes, data: bytes) -> bytes:
        payload = tag + data
        return (_struct.pack(">I", len(data)) + payload
                + _struct.pack(">I", _zlib.crc32(payload) & 0xFFFFFFFF))
    return (
        b"\x89PNG\r\n\x1a\n"
        + _chunk(b"IHDR", _struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
        + _chunk(b"IDAT", _zlib.compress(b"\x00\xFF\xFF\xFF"))
        + _chunk(b"IEND", b"")
    )


_AUDIO_ICON_PNG = _make_tiny_png()

# Audio shape — a <p:pic> element matching exactly what PowerPoint generates when
# you insert audio via Insert > Audio. Critical details verified via COM automation:
#   1. <a:audioFile> uses DRAWINGML namespace (a:), NOT presentationml (p:)
#   2. <p:nvPr> must contain a <p14:media r:embed="{media_rid}"> extension
#      where media_rid is a SECOND relationship with type
#      http://schemas.microsoft.com/office/2007/relationships/media
#   3. <a:blip r:embed> references the icon PNG (image relationship)
_PIC_XML_TMPL = (
    '<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
    ' xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main">'
    '<p:nvPicPr>'
    '<p:cNvPr id="{sid}" name="Audio {idx}">'
    '<a:hlinkClick r:id="" action="ppaction://media"/>'
    '</p:cNvPr>'
    '<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
    '<p:nvPr>'
    '<a:audioFile r:link="{rid}"/>'   # a: namespace — the critical fix
    '<p:extLst>'
    '<p:ext uri="{{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}}">'
    '<p14:media r:embed="{media_rid}"/>'
    '</p:ext>'
    '</p:extLst>'
    '</p:nvPr>'
    '</p:nvPicPr>'
    '<p:blipFill>'
    '<a:blip r:embed="{img_rid}"/>'
    '<a:stretch><a:fillRect/></a:stretch>'
    '</p:blipFill>'
    '<p:spPr>'
    '<a:xfrm><a:off x="457200" y="5943600"/><a:ext cx="457200" cy="457200"/></a:xfrm>'
    '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
    '</p:spPr>'
    '</p:pic>'
)

# Timing template matching what PowerPoint generates for auto-play on slide entry.
# Key differences from naive implementations:
#   - Uses <p:seq concurrent="1" nextAc="seek"> with nodeType="mainSeq"
#   - delay="0" on outer par = auto-triggers without user click
#   - Includes <p:audio> media controller and prev/next condition lists
_TIMING_XML_TMPL = (
    '<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
    '<p:tnLst><p:par>'
    '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
    '<p:childTnLst>'
    '<p:seq concurrent="1" nextAc="seek">'
    '<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
    '<p:childTnLst><p:par>'
    '<p:cTn id="3" fill="hold">'
    '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
    '<p:childTnLst><p:par>'
    '<p:cTn id="4" fill="hold">'
    '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
    '<p:childTnLst><p:par>'
    '<p:cTn id="5" presetID="1" presetClass="mediacall" presetSubtype="0" fill="hold" nodeType="withEffect">'
    '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
    '<p:childTnLst>'
    '<p:cmd type="call" cmd="playFrom(0.0)">'
    '<p:cBhvr>'
    '<p:cTn id="6" dur="indefinite" fill="hold"/>'
    '<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
    '</p:cBhvr>'
    '</p:cmd>'
    '</p:childTnLst>'
    '</p:cTn>'
    '</p:par></p:childTnLst>'
    '</p:cTn>'
    '</p:par></p:childTnLst>'
    '</p:cTn>'
    '</p:par></p:childTnLst>'
    '</p:cTn>'
    '<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
    '<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
    '</p:seq>'
    '<p:audio><p:cMediaNode vol="80000">'
    '<p:cTn id="7" fill="hold" display="0">'
    '<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
    '<p:endCondLst><p:cond evt="onStopAudio" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:endCondLst>'
    '</p:cTn>'
    '<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
    '</p:cMediaNode></p:audio>'
    '</p:childTnLst>'
    '</p:cTn>'
    '</p:par></p:tnLst>'
    '</p:timing>'
)




def embed_audio_into_pptx(
    pptx_bytes: bytes,
    slide_audio: list,          # list[bytes | None], indexed by PPTX slide (0-based)
) -> bytes:
    """
    Embeds MP3 audio into each slide using python-pptx's OPC Part API,
    which correctly manages the package relationships and content-type registry.
    Returns the modified PPTX as bytes.
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    pkg = prs.part.package

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx >= len(slide_audio):
            break
        audio = slide_audio[slide_idx]
        if not audio:
            continue

        # ── add MP3 as a proper OPC Part — TWO relationships needed ──────────────
        # 1. Legacy audio rel (for a:audioFile r:link)
        audio_part = Part(
            PackURI(f"/ppt/media/audio_slide{slide_idx + 1}.mp3"),
            "audio/mpeg",
            pkg,
            audio,
        )
        rid = slide.part.relate_to(audio_part, _AUDIO_REL_TYPE)
        # 2. Modern media rel (for p14:media r:embed) — points to same file
        media_rid = slide.part.relate_to(audio_part, _MEDIA_REL_TYPE)

        # ── add PNG icon as a proper OPC Part ────────────────────────────
        icon_part = Part(
            PackURI(f"/ppt/media/audio_icon_slide{slide_idx + 1}.png"),
            "image/png",
            pkg,
            _AUDIO_ICON_PNG,
        )
        img_rid = slide.part.relate_to(icon_part, _IMAGE_REL_TYPE)

        # ── find next available shape ID ──────────────────────────────────
        sp_tree = slide.shapes._spTree
        existing_ids = [
            int(el.get("id", "0"))
            for el in sp_tree.iter(qn("p:cNvPr"))
        ]
        sid = max(existing_ids, default=0) + 1

        # ── append <p:pic> to spTree ──────────────────────────────────────
        pic_xml = _PIC_XML_TMPL.format(sid=sid, idx=slide_idx + 1, rid=rid, media_rid=media_rid, img_rid=img_rid)
        sp_tree.append(etree.fromstring(pic_xml.encode()))

        # ── replace / insert <p:timing> on the slide element ─────────────
        slide_el = slide._element
        for old in slide_el.findall(qn("p:timing")):
            slide_el.remove(old)
        timing_xml = _TIMING_XML_TMPL.format(sid=sid)
        slide_el.append(etree.fromstring(timing_xml.encode()))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

