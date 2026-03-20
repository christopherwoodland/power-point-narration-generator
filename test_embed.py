"""Quick smoke test for the pptx_builder — no TTS, uses a tiny silent MP3."""
import sys, pathlib
sys.path.insert(0, str(pathlib.Path(__file__).parent / "backend"))
from pptx_builder import embed_audio_into_pptx
from pptx import Presentation
import io

# Minimal valid MP3 (ID3v2 header only - silent, ~100 bytes)
SILENT_MP3 = (
    b"ID3\x03\x00\x00\x00\x00\x00\x23"       # ID3 header
    b"TIT2\x00\x00\x00\x09\x00\x00\x00test"  # TIT2 frame
    b"\xff\xfb\x90\x00" + b"\x00" * 30        # single MP3 frame
)

pptx_bytes = open("test.pptx", "rb").read()
prs_before = Presentation(io.BytesIO(pptx_bytes))
n = len(prs_before.slides)
print(f"PPTX has {n} slides")

# Apply audio to first 2 slides only
audio_list = [SILENT_MP3, SILENT_MP3] + [None] * (n - 2)
result = embed_audio_into_pptx(pptx_bytes, audio_list)
print(f"Output size: {len(result)} bytes")

# Verify result still loads as valid PPTX
prs_after = Presentation(io.BytesIO(result))
print(f"After embed: {len(prs_after.slides)} slides — OK")

with open("test_output.pptx", "wb") as f:
    f.write(result)
print("Wrote test_output.pptx")
